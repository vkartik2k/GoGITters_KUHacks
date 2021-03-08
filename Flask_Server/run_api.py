from flask import Flask
from bot import bot_call
from pptx import Presentation
from flask import jsonify
from flask import request
import requests
# make bot into a function
app = Flask(__name__)

def extract_text(path, marker=None):
    prs = Presentation(path)
    text_runs = []
    if marker is None:
        marker = len(prs.slides)
    c = 0
    for slide in prs.slides:
        c += 1
        if(c > marker):
            break
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    final_text = ""
    for t in text_runs:
        if(len(t.split(" ")) > 4):
            final_text += ("\n"+t)

    return final_text

@app.route('/postAttendance/',methods=["GET","POST"])  # make post 
def attendance():
    print(request.json)
    link=request.json['link']
    # password=request.args.get('password',None)
    password = "USER@123"
    lst =bot_call(link,password)
    return jsonify({"present": lst[1:]})
    
@app.route('/postQuestions/',methods=["GET","POST"]) # PPT and initialization ppt
def getQuestions():
    from app import QuestionGenerator
    # print("Jushaan")
    link=request.json['classNote']
    #get rerquest and download
    response = requests.get(link)
    page=request.json['pageNumber']
    sentences = extract_text('sample.pptx',int(page))
    qg = QuestionGenerator(sentences)
    questions = qg.process(top_sentences=10)

    return_lis =[]
    possible_ans =[]
    for item in questions:
        possible_ans.append(item['answer'])
    possible_ans=set(possible_ans)
    if len(possible_ans)<4:
        print("Need to repeat option???")
    for item in questions:
        questions = item['question_gap']
        ans = item['answer']
        options=[]
        options.append(ans)
        for i,it in enumerate(possible_ans):
            if len(options)==4:
                break
            if it!=ans:
                options.append(it)
        return_lis.append({'correctAns':ans,'statement':questions,'options':options})
    return(jsonify(return_lis)) 
    
# @app.route('/Questions/<page>',methods=["GET"])

if __name__ == '__main__':
    app.run(debug=True)